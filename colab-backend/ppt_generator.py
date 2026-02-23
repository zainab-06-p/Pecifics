"""
PowerPoint File Generator - Enhanced with Canva-Style Professional Themes
Creates .pptx files with animations, headers/footers, and image placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime
from lxml import etree
import os

# Professional Theme Definitions (Canva-inspired)
PROFESSIONAL_THEMES = {
    "modern_blue": {
        "name": "Modern Blue",
        "primary": RGBColor(41, 128, 185),      # Professional blue
        "secondary": RGBColor(52, 152, 219),    # Light blue
        "accent": RGBColor(241, 196, 15),       # Gold accent
        "dark": RGBColor(44, 62, 80),           # Dark slate
        "light": RGBColor(236, 240, 241),       # Light gray
        "background": RGBColor(255, 255, 255),   # White
    },
    "sunset_gradient": {
        "name": "Sunset Gradient",
        "primary": RGBColor(255, 107, 107),     # Coral red
        "secondary": RGBColor(255, 184, 77),    # Orange
        "accent": RGBColor(72, 52, 212),        # Purple
        "dark": RGBColor(46, 64, 83),           # Dark blue
        "light": RGBColor(253, 203, 110),       # Light yellow
        "background": RGBColor(255, 250, 240),   # Cream
    },
    "corporate_green": {
        "name": "Corporate Green",
        "primary": RGBColor(39, 174, 96),       # Green
        "secondary": RGBColor(46, 204, 113),    # Light green
        "accent": RGBColor(230, 126, 34),       # Orange
        "dark": RGBColor(44, 62, 80),           # Slate
        "light": RGBColor(236, 240, 241),       # Light gray
        "background": RGBColor(255, 255, 255),   # White
    },
    "minimalist_dark": {
        "name": "Minimalist Dark",
        "primary": RGBColor(236, 240, 241),     # Light text
        "secondary": RGBColor(149, 165, 166),   # Gray
        "accent": RGBColor(26, 188, 156),       # Teal
        "dark": RGBColor(52, 73, 94),           # Dark blue
        "light": RGBColor(189, 195, 199),       # Silver
        "background": RGBColor(44, 62, 80),      # Dark background
    },
    "vibrant_purple": {
        "name": "Vibrant Purple",
        "primary": RGBColor(155, 89, 182),      # Purple
        "secondary": RGBColor(142, 68, 173),    # Dark purple
        "accent": RGBColor(241, 196, 15),       # Yellow
        "dark": RGBColor(44, 62, 80),           # Slate
        "light": RGBColor(245, 246, 250),       # Off-white
        "background": RGBColor(255, 255, 255),   # White
    },
    "ocean_blue": {
        "name": "Ocean Blue",
        "primary": RGBColor(3, 169, 244),       # Sky blue
        "secondary": RGBColor(0, 188, 212),     # Cyan
        "accent": RGBColor(255, 193, 7),        # Amber
        "dark": RGBColor(38, 50, 56),           # Blue gray
        "light": RGBColor(224, 242, 241),       # Pale cyan
        "background": RGBColor(250, 251, 255),   # Light blue tint
    },
}

def add_slide_transition(slide, transition_type="fade", duration=500):
    """
    Add transition effect to a slide
    
    Args:
        slide: The slide object
        transition_type: Type of transition ('fade', 'push', 'wipe', 'cover', 'uncover')
        duration: Duration in milliseconds (default: 500ms)
    """
    try:
        # Access the slide's XML element
        sld = slide._element
        
        # Create transition element if it doesn't exist
        transition = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        if transition is None:
            transition = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        
        # Set transition speed (duration)
        transition.set('spd', 'med')  # slow, med, fast
        transition.set('advTm', str(duration))
        
        # Add specific transition effect
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        # Clear existing transition effects
        for child in list(transition):
            transition.remove(child)
        
        if transition_type == "fade":
            # Fade transition
            fade = etree.SubElement(transition, f'{p_ns}fade')
            fade.set('thruBlk', '0')
        elif transition_type == "push":
            # Push transition
            push = etree.SubElement(transition, f'{p_ns}push')
            push.set('dir', 'l')  # left
        elif transition_type == "wipe":
            # Wipe transition
            wipe = etree.SubElement(transition, f'{p_ns}wipe')
            wipe.set('dir', 'l')
        elif transition_type == "cover":
            # Cover transition
            cover = etree.SubElement(transition, f'{p_ns}cover')
            cover.set('dir', 'l')
        else:
            # Default to fade
            fade = etree.SubElement(transition, f'{p_ns}fade')
            fade.set('thruBlk', '0')
        
        return True
    except Exception as e:
        print(f"⚠️ Could not add transition: {e}")
        return False

def add_shape_animations(slide, animation_type="fade"):
    """
    Add entrance animations to all shapes on a slide
    
    Args:
        slide: The slide object
        animation_type: Type of animation ('fade', 'fly_in', 'wipe', 'zoom', 'bounce')
    """
    try:
        # Access slide XML
        sld = slide._element
        
        # Namespaces
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        # Get or create timing element
        timing = sld.find(f'.//{p_ns}timing')
        if timing is None:
            timing = etree.SubElement(sld, f'{p_ns}timing')
        
        # Create tnLst (timing node list)
        tnLst = timing.find(f'.//{p_ns}tnLst')
        if tnLst is None:
            tnLst = etree.SubElement(timing, f'{p_ns}tnLst')
        
        # Create parallel time node
        par = etree.SubElement(tnLst, f'{p_ns}par')
        cTn = etree.SubElement(par, f'{p_ns}cTn')
        cTn.set('id', '1')
        cTn.set('dur', 'indefinite')
        cTn.set('restart', 'never')
        cTn.set('nodeType', 'tmRoot')
        
        # Create child time node list
        childTnLst = etree.SubElement(cTn, f'{p_ns}childTnLst')
        
        # Get all shapes to animate
        shapes_to_animate = []
        for idx, shape in enumerate(slide.shapes):
            # Get shape ID from XML
            try:
                shape_id = shape._element.attrib.get('id', str(idx))
                shapes_to_animate.append((shape, shape_id, idx))
            except:
                continue
        
        # Add animation for each shape
        animation_delay = 0
        for shape, shape_id, idx in shapes_to_animate:
            # Create sequence for this shape
            seq = etree.SubElement(childTnLst, f'{p_ns}seq')
            seq.set('concurrent', '1')
            seq.set('nextAc', 'seek')
            
            seqCTn = etree.SubElement(seq, f'{p_ns}cTn')
            seqCTn.set('id', str(idx + 2))
            seqCTn.set('dur', 'indefinite')
            seqCTn.set('nodeType', 'mainSeq')
            
            # Create stCondLst (start condition list)
            stCondLst = etree.SubElement(seqCTn, f'{p_ns}stCondLst')
            cond = etree.SubElement(stCondLst, f'{p_ns}cond')
            cond.set('evt', 'onBegin')
            cond.set('delay', str(animation_delay))
            
            # Create prevCondLst
            prevCondLst = etree.SubElement(seqCTn, f'{p_ns}prevCondLst')
            prevCond = etree.SubElement(prevCondLst, f'{p_ns}cond')
            prevCond.set('evt', 'onPrev')
            prevCond.set('delay', '0')
            
            # Create nextCondLst
            nextCondLst = etree.SubElement(seqCTn, f'{p_ns}nextCondLst')
            nextCond = etree.SubElement(nextCondLst, f'{p_ns}cond')
            nextCond.set('evt', 'onNext')
            nextCond.set('delay', '0')
            
            # Create childTnLst for animation effect
            seqChildTnLst = etree.SubElement(seqCTn, f'{p_ns}childTnLst')
            
            # Create parallel node for effect
            effectPar = etree.SubElement(seqChildTnLst, f'{p_ns}par')
            effectCTn = etree.SubElement(effectPar, f'{p_ns}cTn')
            effectCTn.set('id', str(idx + 100))
            effectCTn.set('fill', 'hold')
            
            # Create start condition
            effectStCondLst = etree.SubElement(effectCTn, f'{p_ns}stCondLst')
            effectCond = etree.SubElement(effectStCondLst, f'{p_ns}cond')
            effectCond.set('delay', '0')
            
            # Create childTnLst for actual animation
            effectChildTnLst = etree.SubElement(effectCTn, f'{p_ns}childTnLst')
            
            # Add animation effect based on type
            if animation_type == "fade":
                # Fade in effect
                set_elem = etree.SubElement(effectChildTnLst, f'{p_ns}set')
                setCTn = etree.SubElement(set_elem, f'{p_ns}cTn')
                setCTn.set('id', str(idx + 200))
                setCTn.set('dur', '1')
                setCTn.set('fill', 'hold')
                
                # Set start condition
                setStCondLst = etree.SubElement(setCTn, f'{p_ns}stCondLst')
                setStartCond = etree.SubElement(setStCondLst, f'{p_ns}cond')
                setStartCond.set('delay', '0')
                
                # Set target
                tgtEl = etree.SubElement(set_elem, f'{p_ns}tgtEl')
                spTgt = etree.SubElement(tgtEl, f'{p_ns}spTgt')
                spTgt.set('spid', str(shape_id))
                
                # Animate effect
                animEffect = etree.SubElement(effectChildTnLst, f'{p_ns}animEffect')
                animEffect.set('transition', 'in')
                animEffect.set('filter', 'fade')
                
                animEffectCTn = etree.SubElement(animEffect, f'{p_ns}cTn')
                animEffectCTn.set('id', str(idx + 300))
                animEffectCTn.set('dur', '500')
                
            elif animation_type == "fly_in":
                # Fly in from left
                animMotion = etree.SubElement(effectChildTnLst, f'{p_ns}animMotion')
                animMotion.set('origin', 'layout')
                animMotion.set('path', 'M 0 0 L 0.25 0')
                animMotion.set('pathEditMode', 'relative')
                
                animMotionCTn = etree.SubElement(animMotion, f'{p_ns}cTn')
                animMotionCTn.set('id', str(idx + 200))
                animMotionCTn.set('dur', '500')
                
                # Target
                tgtEl = etree.SubElement(animMotion, f'{p_ns}tgtEl')
                spTgt = etree.SubElement(tgtEl, f'{p_ns}spTgt')
                spTgt.set('spid', str(shape_id))
            
            # Increment delay for next shape (stagger effect)
            animation_delay += 200  # 200ms delay between each shape
        
        return True
    except Exception as e:
        print(f"⚠️ Could not add shape animations: {e}")
        import traceback
        traceback.print_exc()
        return False

def create_powerpoint(title, slides_data, theme="modern_blue", output_path=None, add_animations=True, add_footer=True):
    """
    Create a PowerPoint presentation with professional Canva-style themes
    
    Args:
        title: Presentation title
        slides_data: List of dicts with 'title' and 'content' keys
        theme: Theme name (modern_blue, sunset_gradient, corporate_green, minimalist_dark, vibrant_purple, ocean_blue)
        output_path: Where to save the file
        add_animations: Add slide transitions and animations (default: True)
        add_footer: Add footer with slide numbers and date (default: True)
    
    Returns:
        Path to created file
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get theme colors
    theme_key = theme.lower().replace(" ", "_")
    if theme_key not in PROFESSIONAL_THEMES:
        theme_key = "modern_blue"  # Default
    colors = PROFESSIONAL_THEMES[theme_key]
    
    # ========== TITLE SLIDE with Professional Design ==========
    title_slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["background"]
    
    # Add large decorative circle (top-right accent)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6.5), Inches(-1),
        Inches(4.5), Inches(4.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors["accent"]
    circle.line.fill.background()
    circle.fill.transparency = 0.3  # 70% opacity
    
    # Add decorative triangle (bottom-left)
    triangle = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(-0.5), Inches(5),
        Inches(3), Inches(3)
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = colors["secondary"]
    triangle.line.fill.background()
    triangle.fill.transparency = 0.4
    
    # Main title text box with shadow effect
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = colors["primary"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle with icon/decoration
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    current_date = datetime.now().strftime("%B %Y")
    subtitle_frame.text = f"Professional Presentation | {current_date}"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = colors["dark"]
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add slide transition for title slide
    if add_animations:
        add_slide_transition(slide, transition_type="fade", duration=800)
        # Add entrance animations to title slide shapes
        add_shape_animations(slide, animation_type="fade")
    
    # ========== CONTENT SLIDES with Professional Design ==========
    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["background"]
        
        # Add header accent bar
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = colors["primary"]
        header_bar.line.fill.background()
        
        # Slide title on accent bar
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_info['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # White on colored bar
        
        # Add image placeholder with icon (right side)
        image_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.5), Inches(1.8),
            Inches(2.2), Inches(2.2)
        )
        image_placeholder.fill.solid()
        image_placeholder.fill.fore_color.rgb = colors["light"]
        image_placeholder.line.color.rgb = colors["secondary"]
        image_placeholder.line.width = Pt(2)
        
        # Add "image" icon text in placeholder
        img_text_box = slide.shapes.add_textbox(Inches(7.5), Inches(2.6), Inches(2.2), Inches(0.6))
        img_frame = img_text_box.text_frame
        img_frame.text = "📷 Image"
        img_para = img_frame.paragraphs[0]
        img_para.font.size = Pt(20)
        img_para.font.color.rgb = colors["secondary"]
        img_para.alignment = PP_ALIGN.CENTER
        
        # Content area with bullet points
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(6.5), Inches(4.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Parse and style content
        content_lines = slide_info['content'].split('\n')
        for i, line in enumerate(content_lines):
            if not line.strip():
                continue
                
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Remove bullet if present, we'll add custom styling
            clean_line = line.strip().lstrip('•').strip()
            
            if line.strip().startswith('•') or '•' in line:
                p.text = f"◆  {clean_line}"
                p.font.color.rgb = colors["dark"]
                p.font.size = Pt(20)
                p.space_before = Pt(8)
            else:
                # Section headers (lines without bullets)
                p.text = clean_line
                p.font.color.rgb = colors["secondary"]
                p.font.size = Pt(22)
                p.font.bold = True
                p.space_before = Pt(12)
        
        # Add decorative corner element (bottom-right)
        corner_accent = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.5), Inches(6.8),
            Inches(1.4), Inches(0.6)
        )
        corner_accent.fill.solid()
        corner_accent.fill.fore_color.rgb = colors["accent"]
        corner_accent.line.fill.background()
        
        # Add footer if requested
        if add_footer:
            # Footer background bar
            footer_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(7.1),
                Inches(10), Inches(0.4)
            )
            footer_bar.fill.solid()
            footer_bar.fill.fore_color.rgb = colors["light"]
            footer_bar.line.fill.background()
            
            # Footer text (left side - presentation title)
            footer_left = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(5), Inches(0.3))
            footer_left_frame = footer_left.text_frame
            footer_left_frame.text = title[:50]  # Truncate if too long
            footer_left_para = footer_left_frame.paragraphs[0]
            footer_left_para.font.size = Pt(10)
            footer_left_para.font.color.rgb = colors["dark"]
            
            # Footer text (right side - slide number)
            footer_right = slide.shapes.add_textbox(Inches(8.5), Inches(7.15), Inches(1.2), Inches(0.3))
            footer_right_frame = footer_right.text_frame
            footer_right_frame.text = f"Slide {idx + 1}/{len(slides_data)}"
            footer_right_para = footer_right_frame.paragraphs[0]
            footer_right_para.font.size = Pt(10)
            footer_right_para.font.color.rgb = colors["dark"]
            footer_right_para.alignment = PP_ALIGN.RIGHT
        
        # Add slide transition (alternating between push and wipe for variety)
        if add_animations:
            if idx % 2 == 0:
                add_slide_transition(slide, transition_type="push", duration=600)
            else:
                add_slide_transition(slide, transition_type="wipe", duration=600)
            
            # Add entrance animations to all shapes on this slide
            add_shape_animations(slide, animation_type="fade")
    
    # Save presentation
    if output_path is None:
        output_path = os.path.join(os.path.expanduser("~"), "Documents", f"{title.replace(' ', '_')}.pptx")
    
    # Ensure directory exists before saving
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
    
    prs.save(output_path)
    return output_path


def generate_slides_from_topic(topic, num_slides=4):
    """Generate slide content based on topic"""
    
    slide_templates = {
        1: {
            "title": "Introduction",
            "content": f"What is {topic}?\n\n• Key characteristics\n• Main features\n• Why it matters\n• Real-world applications"
        },
        2: {
            "title": "Core Concepts",
            "content": "Fundamental Principles:\n\n• Core concept 1\n• Core concept 2\n• Core concept 3\n• Practical applications"
        },
        3: {
            "title": "Advanced Topics",
            "content": "Advanced Features:\n\n• Advanced topic 1\n• Advanced topic 2\n• Advanced topic 3\n• Best practices"
        },
        4: {
            "title": "Conclusion",
            "content": "Key Takeaways:\n\n• Summary point 1\n• Summary point 2\n• Summary point 3\n• Next steps"
        },
        5: {
            "title": "Resources",
            "content": "Learn More:\n\n• Documentation\n• Tutorials\n• Community forums\n• Further reading"
        }
    }
    
    slides = []
    for i in range(1, min(num_slides, 6)):
        if i in slide_templates:
            slides.append(slide_templates[i])
    
    return slides


if __name__ == "__main__":
    # Test
    slides = generate_slides_from_topic("Linux Operating System", 4)
    output = create_powerpoint("Linux OS: A Complete Guide", slides, "ion")
    print(f"Created: {output}")
