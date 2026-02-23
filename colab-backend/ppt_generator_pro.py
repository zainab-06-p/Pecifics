"""
PowerPoint Generator Pro - Gamma/Canva Quality
Advanced layouts, real images, professional animations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime
from lxml import etree
import os
import requests
import io
from PIL import Image

# ===================== GAMMA/CANVA-STYLE THEMES =====================
PREMIUM_THEMES = {
    "gamma_modern": {
        "name": "Gamma Modern",
        "primary": RGBColor(99, 102, 241),      # Indigo
        "secondary": RGBColor(139, 92, 246),    # Purple
        "accent": RGBColor(251, 191, 36),       # Amber
        "dark": RGBColor(17, 24, 39),           # Gray-900
        "light": RGBColor(243, 244, 246),       # Gray-100
        "background": RGBColor(255, 255, 255),  # White
        "gradient_start": RGBColor(99, 102, 241),
        "gradient_end": RGBColor(139, 92, 246),
    },
    "canva_creative": {
        "name": "Canva Creative",
        "primary": RGBColor(0, 201, 167),       # Turquoise
        "secondary": RGBColor(103, 58, 183),    # Deep purple
        "accent": RGBColor(255, 111, 0),        # Bright orange
        "dark": RGBColor(33, 33, 33),           # Almost black
        "light": RGBColor(250, 250, 250),       # Off-white
        "background": RGBColor(255, 255, 255),  # White
        "gradient_start": RGBColor(0, 201, 167),
        "gradient_end": RGBColor(103, 58, 183),
    },
    "minimalist_pro": {
        "name": "Minimalist Pro",
        "primary": RGBColor(0, 0, 0),           # Black
        "secondary": RGBColor(82, 82, 82),      # Gray
        "accent": RGBColor(255, 215, 0),        # Gold
        "dark": RGBColor(0, 0, 0),              # Black
        "light": RGBColor(245, 245, 245),       # Whisper
        "background": RGBColor(255, 255, 255),  # White
        "gradient_start": RGBColor(0, 0, 0),
        "gradient_end": RGBColor(82, 82, 82),
    },
    "sunset_pro": {
        "name": "Sunset Professional",
        "primary": RGBColor(244, 63, 94),       # Rose
        "secondary": RGBColor(251, 146, 60),    # Orange
        "accent": RGBColor(250, 204, 21),       # Yellow
        "dark": RGBColor(127, 29, 29),          # Dark red
        "light": RGBColor(254, 242, 242),       # Rose tint
        "background": RGBColor(255, 251, 235),  # Amber tint
        "gradient_start": RGBColor(244, 63, 94),
        "gradient_end": RGBColor(251, 146, 60),
    },
}

# ===================== UNSPLASH IMAGE INTEGRATION =====================
UNSPLASH_ACCESS_KEY = None  # User will set this

def fetch_unsplash_image(query, orientation="landscape"):
    """Fetch high-quality image from Unsplash API"""
    if not UNSPLASH_ACCESS_KEY:
        return None
    
    try:
        url = "https://api.unsplash.com/photos/random"
        params = {
            "query": query,
            "orientation": orientation,
            "client_id": UNSPLASH_ACCESS_KEY
        }
        response = requests.get(url, params=params, timeout=5)
        if response.status_code == 200:
            data = response.json()
            img_url = data["urls"]["regular"]
            img_response = requests.get(img_url, timeout=10)
            if img_response.status_code == 200:
                return io.BytesIO(img_response.content)
    except Exception as e:
        print(f"⚠️ Image fetch failed: {e}")
    return None

# ===================== ADVANCED SLIDE LAYOUTS =====================

def create_hero_slide(prs, title, subtitle, theme_colors, image_query=None):
    """
    Gamma-style hero slide with large image background and overlay text
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Try to add background image
    if image_query:
        img_stream = fetch_unsplash_image(image_query)
        if img_stream:
            slide.shapes.add_picture(img_stream, 0, 0, 
                                    width=prs.slide_width,
                                    height=prs.slide_height)
    
    # Dark overlay for text readability
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.fill.transparency = 0.5  # 50% dark overlay
    overlay.line.fill.background()
    
    # Large centered title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(66)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(2), Inches(4.2), Inches(6), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content, theme_colors, image_query=None):
    """
    Canva-style two-column layout: Image on left, text on right
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme_colors["background"]
    
    # Header bar with gradient effect
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme_colors["primary"]
    header.line.fill.background()
    
    # Title in header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # LEFT COLUMN - Image
    left_x = Inches(0.5)
    left_y = Inches(1.5)
    left_width = Inches(4.5)
    left_height = Inches(5)
    
    if image_query:
        img_stream = fetch_unsplash_image(image_query, "portrait")
        if img_stream:
            slide.shapes.add_picture(img_stream, left_x, left_y, 
                                    width=left_width, height=left_height)
        else:
            # Placeholder with icon
            img_placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left_x, left_y, left_width, left_height
            )
            img_placeholder.fill.solid()
            img_placeholder.fill.fore_color.rgb = theme_colors["light"]
            img_placeholder.line.color.rgb = theme_colors["secondary"]
            img_placeholder.line.width = Pt(2)
    
    # RIGHT COLUMN - Content
    right_x = Inches(5.2)
    right_y = Inches(1.5)
    right_width = Inches(4.3)
    
    content_box = slide.shapes.add_textbox(right_x, right_y, right_width, Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    # Add content with better formatting
    for item in right_content:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = theme_colors["dark"]
        p.space_before = Pt(12)
        p.level = 0
        p.bullet = True
    
    return slide

def create_comparison_slide(prs, title, left_title, left_items, right_title, right_items, theme_colors):
    """
    Gamma-style comparison slide: Before/After, Pros/Cons, etc.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme_colors["background"]
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_colors["primary"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # LEFT COLUMN
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.7)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = theme_colors["primary"]
    left_header.line.fill.background()
    
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.7))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.font.size = Pt(28)
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = RGBColor(255, 255, 255)
    left_title_para.alignment = PP_ALIGN.CENTER
    
    # Left content
    left_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(4.3), Inches(4))
    left_text_frame = left_content_box.text_frame
    for item in left_items:
        p = left_text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = theme_colors["dark"]
        p.space_before = Pt(10)
        p.bullet = True
    
    # RIGHT COLUMN
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.7)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = theme_colors["secondary"]
    right_header.line.fill.background()
    
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.7))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.font.size = Pt(28)
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = RGBColor(255, 255, 255)
    right_title_para.alignment = PP_ALIGN.CENTER
    
    # Right content
    right_content_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.4), Inches(4.3), Inches(4))
    right_text_frame = right_content_box.text_frame
    for item in right_items:
        p = right_text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = theme_colors["dark"]
        p.space_before = Pt(10)
        p.bullet = True
    
    return slide

def create_big_number_slide(prs, title, number, description, theme_colors):
    """
    Gamma-style slide with huge number/stat as focal point
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme_colors["background"]
    
    # Decorative accent shape
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(0.5), Inches(3), Inches(6.5)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = theme_colors["accent"]
    accent_shape.fill.transparency = 0.15
    accent_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(42)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_colors["primary"]
    
    # HUGE number
    number_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(2.5))
    number_frame = number_box.text_frame
    number_frame.text = str(number)
    number_para = number_frame.paragraphs[0]
    number_para.font.size = Pt(120)
    number_para.font.bold = True
    number_para.font.color.rgb = theme_colors["primary"]
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(6), Inches(2))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.word_wrap = True
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(24)
    desc_para.font.color.rgb = theme_colors["dark"]
    
    return slide

def create_quote_slide(prs, quote, author, theme_colors):
    """
    Canva-style quote slide with elegant typography
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background with subtle color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme_colors["light"]
    
    # Large quotation mark decoration
    quote_mark = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(2))
    quote_frame = quote_mark.text_frame
    quote_frame.text = '"'
    quote_para = quote_frame.paragraphs[0]
    quote_para.font.size = Pt(180)
    quote_para.font.color.rgb = theme_colors["accent"]
    quote_para.font.color.rgb.transparency = 0.3
    
    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(3))
    quote_text_frame = quote_box.text_frame
    quote_text_frame.text = quote
    quote_text_frame.word_wrap = True
    quote_text_para = quote_text_frame.paragraphs[0]
    quote_text_para.font.size = Pt(32)
    quote_text_para.font.italic = True
    quote_text_para.font.color.rgb = theme_colors["dark"]
    quote_text_para.alignment = PP_ALIGN.CENTER
    
    # Author
    author_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(7), Inches(0.6))
    author_frame = author_box.text_frame
    author_frame.text = f"— {author}"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(24)
    author_para.font.bold = True
    author_para.font.color.rgb = theme_colors["primary"]
    author_para.alignment = PP_ALIGN.CENTER
    
    return slide

# ===================== ADVANCED ANIMATIONS =====================

def add_zoom_animation(slide, shape_id):
    """Add zoom entrance animation (Gamma-style)"""
    try:
        sld = slide._element
        timing = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
        if timing is None:
            timing = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
        
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        # Create animation node for zoom effect
        tnLst = etree.SubElement(timing, f'{p_ns}tnLst')
        par = etree.SubElement(tnLst, f'{p_ns}par')
        cTn = etree.SubElement(par, f'{p_ns}cTn')
        cTn.set('id', '1')
        cTn.set('dur', 'indefinite')
        
        # Animation sequence
        childTnLst = etree.SubElement(cTn, f'{p_ns}childTnLst')
        seq = etree.SubElement(childTnLst, f'{p_ns}seq')
        seqCTn = etree.SubElement(seq, f'{p_ns}cTn')
        seqCTn.set('id', '2')
        seqCTn.set('dur', 'indefinite')
        
        # Zoom animation effect
        seqChildTnLst = etree.SubElement(seqCTn, f'{p_ns}childTnLst')
        animScale = etree.SubElement(seqChildTnLst, f'{p_ns}animScale')
        
        animScaleCTn = etree.SubElement(animScale, f'{p_ns}cTn')
        animScaleCTn.set('id', '3')
        animScaleCTn.set('dur', '500')
        
        # Target shape
        tgtEl = etree.SubElement(animScale, f'{p_ns}tgtEl')
        spTgt = etree.SubElement(tgtEl, f'{p_ns}spTgt')
        spTgt.set('spid', str(shape_id))
        
        # Scale from 0 to 100%
        byPos = etree.SubElement(animScale, f'{p_ns}by')
        byPos.set('x', '100000')
        byPos.set('y', '100000')
        
        return True
    except Exception as e:
        print(f"⚠️ Zoom animation failed: {e}")
        return False

def add_slide_transition_advanced(slide, transition_type="morph"):
    """
    Advanced slide transitions like Gamma/PowerPoint Designer
    """
    try:
        sld = slide._element
        transition = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        if transition is None:
            transition = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        
        transition.set('spd', 'med')
        
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        # Clear existing
        for child in list(transition):
            transition.remove(child)
        
        if transition_type == "morph":
            # Morph transition (PowerPoint 2019+)
            morph = etree.SubElement(transition, f'{p_ns}morph')
        elif transition_type == "zoom":
            # Zoom transition
            zoom = etree.SubElement(transition, f'{p_ns}zoom')
            zoom.set('dir', 'in')
        elif transition_type == "reveal":
            # Reveal transition
            reveal = etree.SubElement(transition, f'{p_ns}reveal')
            reveal.set('dir', 'l')
        
        return True
    except Exception as e:
        print(f"⚠️ Advanced transition failed: {e}")
        return False

# ===================== MAIN GENERATION FUNCTION =====================

def generate_gamma_presentation(topic, num_slides=5, theme="gamma_modern", unsplash_key=None):
    """
    Generate Gamma/Canva quality presentation with AI content and real images
    """
    global UNSPLASH_ACCESS_KEY
    UNSPLASH_ACCESS_KEY = unsplash_key
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get theme
    theme_key = theme.lower().replace(" ", "_")
    if theme_key not in PREMIUM_THEMES:
        theme_key = "gamma_modern"
    colors = PREMIUM_THEMES[theme_key]
    
    print(f"🎨 Creating {colors['name']} presentation about: {topic}")
    
    # SLIDE 1: Hero slide with background image
    print("  📸 Creating hero slide...")
    hero_slide = create_hero_slide(
        prs, 
        title=topic.title(),
        subtitle=f"Professional Presentation | {datetime.now().strftime('%B %Y')}",
        theme_colors=colors,
        image_query=topic
    )
    add_slide_transition_advanced(hero_slide, "zoom")
    
    # SLIDE 2: Two-column intro slide
    print("  📋 Creating introduction...")
    intro_slide = create_two_column_slide(
        prs,
        title="Overview",
        left_content=[],
        right_content=[
            "Comprehensive exploration",
            "Key concepts and insights",
            "Practical applications",
            "Future perspectives"
        ],
        theme_colors=colors,
        image_query=f"{topic} overview"
    )
    add_slide_transition_advanced(intro_slide, "reveal")
    
    # SLIDE 3: Big number/stat slide
    print("  📊 Creating key metric...")
    stat_slide = create_big_number_slide(
        prs,
        title="Key Impact",
        number="100%",
        description=f"Excellence in {topic.title()} - transforming the way we work",
        theme_colors=colors
    )
    add_slide_transition_advanced(stat_slide, "morph")
    
    # SLIDE 4: Comparison slide
    print("  ⚖️ Creating comparison...")
    comparison_slide = create_comparison_slide(
        prs,
        title="Traditional vs Modern Approach",
        left_title="Traditional",
        left_items=[
            "Manual processes",
            "Limited scalability",
            "Higher costs",
            "Slower results"
        ],
        right_title="Modern",
        right_items=[
            "Automated workflows",
            "Infinite scalability",
            "Cost-effective",
            "Real-time insights"
        ],
        theme_colors=colors
    )
    add_slide_transition_advanced(comparison_slide, "reveal")
    
    # SLIDE 5: Quote/closing slide
    if num_slides >= 5:
        print("  💬 Creating closing slide...")
        quote_slide = create_quote_slide(
            prs,
            quote=f"The future of {topic.lower()} is not just about technology, it's about transforming possibilities into reality.",
            author="Industry Expert",
            theme_colors=colors
        )
        add_slide_transition_advanced(quote_slide, "fade")
    
    print("✅ Gamma-quality presentation created!")
    return prs

def save_presentation(prs, output_path):
    """Save with directory auto-creation"""
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
    
    prs.save(output_path)
    return output_path

# ===================== EXAMPLE USAGE =====================
if __name__ == "__main__":
    # Example: Create Gamma-style presentation
    prs = generate_gamma_presentation(
        topic="Artificial Intelligence",
        num_slides=5,
        theme="gamma_modern",
        unsplash_key=None  # Set to your Unsplash API key for real images
    )
    
    output_file = save_presentation(prs, "AI_Presentation_Pro.pptx")
    print(f"\n📥 Saved to: {output_file}")
    print(f"💡 Tip: Get free Unsplash API key from https://unsplash.com/developers")
