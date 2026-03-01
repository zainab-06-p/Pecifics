"""
Pecifics AI Desktop Assistant — LangGraph Backend
=====================================================
Architecture:
  Brain 1: Groq Llama-3.3-70b  → Plans tasks, splits multi-step prompts
  Brain 2: CogAgent (Kaggle)    → Vision agent: screenshot → coordinates → actions
  
Quick start:
  1. Set GROQ_API_KEY in .env (free: https://console.groq.com)
  2. Set COGAGENT_URL in .env (from Kaggle notebook ngrok URL)
  3. python langchain_backend.py
"""

import os, re, json, base64, logging, traceback, gc, httpx
from io import BytesIO
from typing import Dict, List, Optional
from datetime import datetime

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
import uvicorn

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

# ─── ENV ─────────────────────────────────────────────────────────────────────
try:
    from dotenv import load_dotenv; load_dotenv()
except ImportError:
    pass

GROQ_API_KEY    = os.getenv("GROQ_API_KEY", "")
GROQ_MODEL      = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")
COGAGENT_URL    = os.getenv("COGAGENT_URL", "")
GEMINI_API_KEY  = os.getenv("GEMINI_API_KEY", "")
GEMINI_MODEL    = os.getenv("GEMINI_MODEL", "gemini-2.0-flash")
MAX_RETRIES     = int(os.getenv("MAX_RETRIES", "3"))

import pathlib as _pathlib
_USER_HOME = str(_pathlib.Path.home()).replace("\\", "\\\\")
_USER_NAME = _pathlib.Path.home().name

# ─── DEPS ────────────────────────────────────────────────────────────────────
HAS_GROQ = False
HAS_GEMINI = False

try:
    from langchain_groq import ChatGroq; HAS_GROQ = True
except ImportError: pass

try:
    import google.generativeai as genai; HAS_GEMINI = True
except ImportError: pass

from langchain_core.messages import HumanMessage

# ─── REQUEST MODELS ──────────────────────────────────────────────────────────

class ChatRequest(BaseModel):
    message:              str
    screenshot:           Optional[str] = None
    conversation_history: Optional[List[Dict]] = []
    screen_width:         Optional[int] = 1920
    screen_height:        Optional[int] = 1080
    user_choice:          Optional[Dict] = None
    session_id:           Optional[str] = None
    user_home:            Optional[str] = None

class VisionActRequest(BaseModel):
    screenshot:        str
    goal:              str
    step_history:      List[Dict] = []
    screen_width:      Optional[int] = 1920
    screen_height:     Optional[int] = 1080
    cogagent_url:      Optional[str] = None

class VerifyRequest(BaseModel):
    screenshot:      str
    task:            str
    expected_result: Optional[str] = ""

class NextStepRequest(BaseModel):
    original_task:     str
    last_action:       Dict
    last_result:       Dict
    screenshot:        Optional[str] = None
    remaining_actions: List[Dict] = []
    completed_actions: List[Dict] = []

class GeneratePPTRequest(BaseModel):
    topic:                   str
    title:                   Optional[str] = None
    num_slides:              int = 5
    theme:                   str = "gamma_modern"
    save_path:               str = "Desktop"
    additional_instructions: Optional[str] = None

class AnalyzeScreenRequest(BaseModel):
    screenshot: Optional[str] = None
    question:   Optional[str] = "What do you see on screen?"

class BrowserStateRequest(BaseModel):
    screenshot: str = ""

# ─── ACTION CATALOG ──────────────────────────────────────────────────────────

ACTION_CATALOG = """
=== VISION AGENT (★ USE FOR ALL GUI/BROWSER TASKS) ===
vision_task(goal, max_steps)
    ★ ALWAYS use for ANY browser, website, or GUI interaction.
    AI takes screenshots → identifies UI elements → clicks/types at pixel coordinates.
    goal = detailed description with ALL details (URLs, emails, text)
    max_steps = limit (default 30)

=== SCREEN ANALYSIS ===
get_screenshot(question) | describe_screen() | analyze_screen(question)

=== FILE OPERATIONS ===
create_file(file_path, content) | create_folder(folder_path) | delete_file(file_path)
read_file(file_path) | append_to_file(file_path, content)
find_in_file(file_path, search_text) | replace_in_file(file_path, search_text, replacement_text)
list_directory(path) | copy_file(source, destination)
move_file(source, destination) | rename_file(old_path, new_name)
search_files(pattern, location) | open_file(file_path) | show_in_explorer(file_path)

=== APPLICATION CONTROL ===
open_application(app_name) | close_application(app_name) | open_url(url)
search_web(query)                → Opens Google search results for the query
focus_app(app_name)             → Brings an already-open app to the foreground

=== SYSTEM CONTROLS ===
set_volume(volume)              → 0-100
set_brightness(brightness)      → 0-100 (laptops only)
toggle_wifi(enable)             → enable=true/false, needs admin for some adapters
toggle_bluetooth(enable)        → enable=true/false, needs admin
toggle_night_light(enable)      → enable=true/false (blue light filter)
toggle_dark_mode(enable)        → enable=true/false (system-wide dark mode)
set_wallpaper(image_path)       → Full path to image file
lock_computer() | sleep_computer()
get_battery_status() | get_system_info() | get_disk_space()
get_network_status() | empty_recycle_bin()
speak(message)                  → Text-to-speech
show_notification(title, message) → Windows toast notification
get_clipboard() | set_clipboard(text)

=== OS TASKS ===
clear_cache(cache_type)         → type: 'all', 'temp', 'browser', 'dns'
flush_dns() | reset_network()
open_task_manager() | open_control_panel() | open_device_manager()
manage_service(service_name, action) → action: 'start', 'stop', 'restart'
get_running_processes() | kill_process(name_or_pid)
get_system_health() | get_power_plan() | set_power_plan(plan)
restart_computer(delay) | shutdown_computer(delay) | cancel_shutdown()

=== MS OFFICE — POWERPOINT (AUTO-OPENS ON SCREEN WHEN DONE) ===
generate_ppt(topic, title, num_slides, theme, save_path, additional_instructions)
    → Creates professional .pptx and AUTO-OPENS it in PowerPoint on screen
ppt_find_slide(search_text) | ppt_get_slide_content(slide_number)
ppt_update_slide_text(slide_number, old_text, new_text)
ppt_apply_theme(theme_name) | ppt_add_animation(animation_type)
ppt_change_layout(layout_name)

=== MS OFFICE — WORD (AUTO-OPENS ON SCREEN WHEN DONE) ===
word_create_document(title, content)
    → Creates .docx document and opens it in Word on screen
word_open_document(filepath) | word_read_content()
word_add_paragraph(text) | word_add_heading(text, level)
word_find_replace(search_text, replacement_text)
word_insert_table(rows, cols) | word_apply_theme(theme_name)
word_save(filename) | word_change_font(font_name, font_size)

=== MS OFFICE — EXCEL (AUTO-OPENS ON SCREEN WHEN DONE) ===
excel_create_workbook()
    → Creates .xlsx workbook and opens it in Excel on screen
excel_open_workbook(filepath) | excel_write_cell(row, col, value)
excel_write_data(data) | excel_add_worksheet(name)
excel_create_chart(chart_type) | excel_format_cell(row, col, options)
excel_autofit_columns() | excel_save(filename)

=== MS OFFICE — ONENOTE ===
onenote_open()                  → Opens OneNote application
onenote_create_page(title)      → Creates a new page with given title
onenote_add_content(content)    → Adds text content to the current page
onenote_list_notebooks()        → Lists all notebooks and their sections

=== MS OFFICE — PUBLISHER ===
publisher_create(template_type) → Creates a new publication (blank by default)
publisher_add_textbox(text, left, top, width, height) → Adds a text box
publisher_add_page()            → Adds a new page
publisher_save(filename)        → Saves the publication
"""

# ─── SYSTEM PROMPT ───────────────────────────────────────────────────────────

SYSTEM_PROMPT = f"""You are Pecifics, an AI desktop assistant for Windows 11.
You work like Claude Computer Use — you see the screen and control it.

EXECUTION MODEL:
- For ALL browser/GUI tasks → vision_task(goal="..."). The vision agent sees the screen and clicks.
- For file/system/non-GUI tasks → use the specific action directly.
- NEVER use browser_click, browser_type, or CSS selectors.

MULTI-TASK HANDLING:
- Users may give 3-4 tasks in one prompt.
- You MUST split them and return ALL tasks in the tasks array.
- Example: "Open Chrome, create a file on desktop, and set volume to 50"
  → 3 separate tasks, each with its own actions.

USER INFO:
- Home: {_USER_HOME}  |  Username: {_USER_NAME}
- ALWAYS use real path. NEVER use "USERNAME" placeholder.

ASKING FOR INPUT:
- If task is missing CRITICAL info → set needs_input=true with input_fields
- Needs input ONLY WHEN NOT PROVIDED: save file (filename+location), send email (to address if missing),
  create PPT (topic if unclear), delete (confirmation needed)
- If user already provided enough info (e.g. "write a mail saying I won't attend tomorrow"),
  DO NOT ask for input — use vision_task with a detailed goal instead.
- No input needed: open apps, volume, system info, open websites, compose email with given content

{ACTION_CATALOG}

OUTPUT FORMAT (ALWAYS valid JSON, no markdown):
{{
  "message": "Brief summary of all tasks",
  "tasks": [
    {{
      "id": 1,
      "description": "Open Chrome and search for weather",
      "needs_input": false,
      "input_fields": [],
      "actions": [{{"name": "vision_task", "parameters": {{"goal": "Open Chrome. Go to google.com. Type 'weather'. Press Enter."}}}}]
    }},
    {{
      "id": 2,
      "description": "Create a text file",
      "needs_input": true,
      "input_fields": [
        {{"key": "filename", "label": "File Name", "placeholder": "e.g. notes.txt", "default": "notes.txt"}},
        {{"key": "content", "label": "Content", "placeholder": "What to write?", "default": ""}}
      ],
      "actions": []
    }}
  ],
  "expected_result": "Chrome shows weather results and file is created"
}}

VISION TASK GOALS — be VERY detailed:
  "Open Chrome. Navigate to https://mail.google.com. Click Compose button. Type bob@test.com in To field. Type 'Meeting' in Subject. Type email body. Click Send."

KEY RULES:
1. Include ALL info in vision_task goals (URLs, emails, text content)
2. Start goals with "Open Chrome." if browser isn't open
3. Gmail user is already logged in — skip login steps
4. For PPT (any "create presentation/slides/ppt"): use generate_ppt action → it AUTO-OPENS in PowerPoint
5. For Word docs ("create document/report/letter"): use word_create_document → AUTO-OPENS in Word
6. For Excel ("create spreadsheet/workbook/table"): use excel_create_workbook → AUTO-OPENS in Excel
7. For file ops: use create_file, read_file etc. (not vision_task)
8. NEVER use generate_ppt + vision_task for same PPT — generate_ppt does everything
9. ALWAYS return tasks array even for single tasks
10. ALL Office actions open on screen automatically — no extra open_file needed"""

# ─── LLM ─────────────────────────────────────────────────────────────────────

_llm = None

def create_llm():
    if not HAS_GROQ or not GROQ_API_KEY:
        raise ValueError("GROQ_API_KEY not set. Free key: https://console.groq.com")
    return ChatGroq(api_key=GROQ_API_KEY, model_name=GROQ_MODEL, temperature=0.1)

def get_llm():
    global _llm
    if _llm is None:
        _llm = create_llm()
    return _llm

def _call(llm, prompt: str) -> str:
    try:
        resp = llm.invoke([HumanMessage(content=prompt)])
        text = resp.content if hasattr(resp, "content") else str(resp)
        logger.debug(f"LLM raw ({len(text)} chars): {text[:300]}")
        return text
    except TypeError:
        return llm.invoke(prompt)
    except Exception as e:
        logger.error(f"LLM call failed: {e}")
        raise

# ─── VISION PROVIDER ─────────────────────────────────────────────────────────

async def call_vision_provider(screenshot_b64: str, goal: str, step_history: list,
                                w: int, h: int, cogagent_url: str = None) -> dict:
    """Call CogAgent (Kaggle) or Gemini for vision actions."""
    
    # Priority 1: CogAgent on Kaggle (prefer dynamic URL from frontend)
    effective_url = cogagent_url or COGAGENT_URL
    if effective_url:
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                resp = await client.post(f"{effective_url.rstrip('/')}/vision_act", json={
                    "screenshot": screenshot_b64, "goal": goal,
                    "step_history": step_history,
                    "screen_width": w, "screen_height": h,
                })
                if resp.status_code == 200:
                    result = resp.json()
                    result.setdefault("action", "fail")
                    result.setdefault("description", "")
                    logger.info(f"CogAgent → {result.get('action')}: {result.get('description','')[:80]}")
                    return result
                logger.warning(f"CogAgent {resp.status_code}, falling back")
        except Exception as e:
            logger.warning(f"CogAgent error: {e}, falling back")
    
    # Priority 2: Gemini
    if HAS_GEMINI and GEMINI_API_KEY:
        return await _gemini_vision_act(screenshot_b64, goal, step_history, w, h)
    
    return {"action": "fail", "description": "No vision provider. Set COGAGENT_URL or GEMINI_API_KEY."}


VISION_ACT_PROMPT = """You are a SCREEN AGENT controlling a computer by looking at screenshots.
Decide the ONE next action to progress the user's goal.

TASK: {goal}

STEPS DONE:
{history}

SCREENSHOT: {w}×{h} pixels. (0,0)=top-left.

RULES:
1. Return EXACTLY ONE action as JSON.
2. Aim for CENTER of UI elements.
3. Dismiss popups/banners FIRST.
4. Click field first, then type next step.
5. Goal achieved → {{"action":"done"}}.
6. Cannot proceed → {{"action":"fail","description":"reason"}}.

ACTIONS:
  {{"action":"click","x":<int>,"y":<int>,"description":"what"}}
  {{"action":"double_click","x":<int>,"y":<int>,"description":"what"}}
  {{"action":"right_click","x":<int>,"y":<int>,"description":"what"}}
  {{"action":"type","text":"...","description":"typing"}}
  {{"action":"type","x":<int>,"y":<int>,"text":"...","description":"click+type"}}
  {{"action":"key","key":"Enter","description":"pressing Enter"}}
  {{"action":"scroll","direction":"down","clicks":3,"description":"scrolling"}}
  {{"action":"wait","duration":2000,"description":"waiting"}}
  {{"action":"done","description":"complete"}}
  {{"action":"fail","description":"cannot proceed"}}

Return ONLY the JSON object."""


async def _gemini_vision_act(b64, goal, history, w, h):
    history_str = "\n".join(
        f"  {i+1}. [{s.get('action','?')}] {s.get('description','')}{' ✓' if s.get('success') else ' ✗'}"
        for i, s in enumerate(history)
    ) or "(none)"

    prompt = VISION_ACT_PROMPT.format(goal=goal, history=history_str, w=w, h=h)
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel(GEMINI_MODEL)
        img_data = base64.b64decode(b64)
        resp = model.generate_content(
            [prompt, {"mime_type": "image/jpeg", "data": img_data}],
            generation_config={"temperature": 0.1, "max_output_tokens": 512}
        )
        raw = resp.text.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        result = json.loads(raw)
        result.setdefault("action", "fail")
        result.setdefault("description", "")
        return result
    except json.JSONDecodeError:
        m = re.search(r'\{[^}]+\}', raw)
        if m:
            try: return json.loads(m.group())
            except: pass
        return {"action": "fail", "description": f"Vision JSON error: {raw[:100]}"}
    except Exception as e:
        err_str = str(e)
        if "429" in err_str or "quota" in err_str.lower():
            logger.error("Gemini API quota exceeded — set COGAGENT_URL to use CogAgent instead")
            return {"action": "fail", "description": "Gemini quota exceeded. Please set your CogAgent URL in Settings (from Kaggle notebook ngrok URL)."}
        return {"action": "fail", "description": f"Vision error: {err_str[:200]}"}


def describe_screenshot(b64: str) -> str:
    if not b64: return "No screenshot."
    if HAS_GEMINI and GEMINI_API_KEY:
        try:
            genai.configure(api_key=GEMINI_API_KEY)
            model = genai.GenerativeModel(GEMINI_MODEL)
            resp = model.generate_content([
                "Describe this Windows screenshot in ≤80 words.",
                {"mime_type": "image/jpeg", "data": base64.b64decode(b64)},
            ])
            return resp.text.strip()
        except: pass
    return "Screenshot captured."

# ─── TASK PLANNER ─────────────────────────────────────────────────────────────

def plan_tasks(user_message, screenshot_b64, conversation_history,
               screen_width, screen_height, user_choice, extra_context="",
               retry_count=0, previous_error=None):
    llm = get_llm()
    
    history_txt = "\n".join(
        f"{'User' if t.get('role')=='user' else 'AI'}: {t.get('content','')}"
        for t in conversation_history[-6:]
    )
    screen_desc = describe_screenshot(screenshot_b64) if screenshot_b64 else "No screenshot."
    retry_ctx = f"\n\nPREVIOUS FAILED: {previous_error}\nUse DIFFERENT approach." if retry_count > 0 and previous_error else ""
    
    if user_choice:
        user_message += f"\n[User selected: {user_choice.get('type')}={user_choice.get('value')}]"

    prompt = f"""{SYSTEM_PROMPT}
{extra_context}
SCREEN: {screen_desc}
SCREEN SIZE: {screen_width}x{screen_height}
HISTORY:
{history_txt}
{retry_ctx}

USER: {user_message}

Return ONLY valid JSON."""

    try:
        raw = _call(llm, prompt).strip()
        logger.info(f"LLM response ({len(raw)} chars): {raw[:500]}")
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        
        # Try direct parse first
        try:
            result = json.loads(raw)
        except json.JSONDecodeError:
            # Extract the outermost JSON object from mixed text
            depth = 0; start = -1; result = None
            for i, c in enumerate(raw):
                if c == '{':
                    if depth == 0: start = i
                    depth += 1
                elif c == '}':
                    depth -= 1
                    if depth == 0 and start >= 0:
                        try:
                            result = json.loads(raw[start:i+1])
                            break
                        except json.JSONDecodeError:
                            start = -1
            if result is None:
                raise json.JSONDecodeError("No valid JSON object found", raw, 0)
        
        result.setdefault("message", "Processing…")
        result.setdefault("expected_result", "")
        
        # Normalize: support old format (actions) → new format (tasks)
        if "tasks" not in result and "actions" in result:
            actions = result.pop("actions", [])
            needs_input = result.get("clarification_needed", False)
            input_fields = result.get("clarification_fields", [])
            result["tasks"] = [{
                "id": 1,
                "description": result.get("task_summary", result["message"]),
                "needs_input": needs_input,
                "input_fields": input_fields,
                "actions": actions if not needs_input else [],
            }]
        
        result.setdefault("tasks", [])
        for i, task in enumerate(result["tasks"]):
            task.setdefault("id", i + 1)
            task.setdefault("description", "")
            task.setdefault("needs_input", False)
            task.setdefault("input_fields", [])
            task.setdefault("actions", [])
        
        return result
    
    except json.JSONDecodeError as e:
        logger.error(f"Plan JSON error (attempt {retry_count+1}): {e}\nRaw: {raw[:500] if 'raw' in dir() else 'N/A'}")
        if retry_count < MAX_RETRIES:
            return plan_tasks(user_message, screenshot_b64, conversation_history,
                             screen_width, screen_height, user_choice, extra_context,
                             retry_count+1, f"JSON error: {e}. Return ONLY valid JSON, no markdown, no extra text.")
        return {"message": "Planning failed. Please rephrase.", "tasks": []}
    except Exception as e:
        logger.error(f"Plan error (attempt {retry_count+1}): {e}")
        if retry_count < MAX_RETRIES:
            return plan_tasks(user_message, screenshot_b64, conversation_history,
                             screen_width, screen_height, user_choice, extra_context,
                             retry_count+1, str(e))
        return {"message": f"Error: {e}", "tasks": []}


def verify_completion(screenshot_b64, task, expected):
    llm = get_llm()
    screen = describe_screenshot(screenshot_b64) if screenshot_b64 else "No screenshot."
    prompt = f"""Task: "{task}"
Expected: "{expected}"
Screen: {screen}

Was it completed? Return ONLY JSON:
{{"success": true, "observation": "what I see", "should_retry": false, "retry_actions": []}}"""
    try:
        raw = _call(llm, prompt).strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        return json.loads(raw)
    except:
        return {"success": False, "observation": "Verification failed", "should_retry": False}


# ─── PPT HELPER ──────────────────────────────────────────────────────────────

def _resolve_save_dir(friendly):
    import pathlib
    home = str(pathlib.Path.home())
    low = friendly.strip().lower()
    if low == "desktop":    return os.path.join(home, "Desktop")
    if low == "documents":  return os.path.join(home, "Documents")
    if low == "downloads":  return os.path.join(home, "Downloads")
    if os.path.isabs(friendly): return friendly
    return os.path.join(home, "Desktop")

# ─── FASTAPI ─────────────────────────────────────────────────────────────────

app = FastAPI(title="Pecifics AI Backend", version="4.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

@app.get("/health")
async def health():
    status = "ok"
    try: get_llm()
    except Exception as e: status = f"error: {e}"
    vision = "cogagent" if COGAGENT_URL else ("gemini" if GEMINI_API_KEY else "none")
    return {"status": "ok", "version": "4.0.0", "llm": f"groq/{GROQ_MODEL}",
            "vision": vision, "cogagent_url": COGAGENT_URL or "not set",
            "timestamp": datetime.now().isoformat()}

@app.post("/chat")
async def chat(req: ChatRequest):
    try: get_llm()
    except Exception as e: raise HTTPException(503, f"LLM unavailable: {e}")
    try:
        extra = ""
        if req.user_home:
            extra = f"\nFRONTEND_USER_HOME: {req.user_home.replace(chr(92), chr(92)*2)}"
        result = plan_tasks(req.message, req.screenshot, req.conversation_history or [],
                           req.screen_width or 1920, req.screen_height or 1080,
                           req.user_choice, extra)
        return JSONResponse(content=result)
    except Exception as e:
        logger.error(traceback.format_exc())
        raise HTTPException(500, str(e))

@app.post("/vision_act")
async def vision_act(req: VisionActRequest):
    # Update global COGAGENT_URL if frontend sent one (persists for session)
    global COGAGENT_URL
    if req.cogagent_url and req.cogagent_url.strip():
        COGAGENT_URL = req.cogagent_url.strip()
        logger.info(f"CogAgent URL updated from request: {COGAGENT_URL}")
    return await call_vision_provider(req.screenshot, req.goal, req.step_history,
                                      req.screen_width or 1920, req.screen_height or 1080,
                                      cogagent_url=req.cogagent_url)

@app.post("/set_config")
async def set_config(req: dict):
    """Update runtime config from frontend settings panel."""
    global COGAGENT_URL, GEMINI_API_KEY
    updated = []
    if req.get("cogagent_url"):
        COGAGENT_URL = req["cogagent_url"].strip()
        updated.append(f"COGAGENT_URL={COGAGENT_URL}")
    if req.get("gemini_api_key"):
        GEMINI_API_KEY = req["gemini_api_key"].strip()
        updated.append("GEMINI_API_KEY=***")
    logger.info(f"Config updated: {updated}")
    return {"success": True, "updated": updated}

@app.post("/verify")
async def verify(req: VerifyRequest):
    return JSONResponse(content=verify_completion(req.screenshot, req.task, req.expected_result))

@app.post("/analyze_screen")
async def analyze_screen(req: AnalyzeScreenRequest):
    desc = describe_screenshot(req.screenshot or "")
    return {"answer": desc, "description": desc}

@app.post("/check_browser_state")
async def check_browser_state(req: BrowserStateRequest):
    if not req.screenshot or not (HAS_GEMINI and GEMINI_API_KEY):
        return {"state": "clear", "can_auto_handle": False, "needs_user": False}
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel(GEMINI_MODEL)
        resp = model.generate_content([
            'Analyze browser screenshot for blockers. Return JSON: {"state":"clear|cookie_consent|2fa|captcha|error","description":"...","can_auto_handle":true/false,"needs_user":true/false,"user_message":"..."}',
            {"mime_type": "image/jpeg", "data": base64.b64decode(req.screenshot)}
        ])
        raw = resp.text.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"): raw = raw[4:]
        return json.loads(raw.strip())
    except:
        return {"state": "clear", "can_auto_handle": False, "needs_user": False}

@app.post("/next_step")
async def next_step(req: NextStepRequest):
    try: llm = get_llm()
    except: return {"decision": "continue", "next_actions": []}
    
    screen = describe_screenshot(req.screenshot) if req.screenshot else "No screenshot."
    prompt = f"""TASK: {req.original_task}
LAST: {json.dumps(req.last_action)} → {json.dumps(req.last_result)}
SCREEN: {screen}
REMAINING: {json.dumps(req.remaining_actions[:3])}

Decide: continue|replace|done. Return JSON: {{"decision":"...","message":"...","next_actions":[]}}"""
    try:
        raw = _call(llm, prompt).strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        result = json.loads(raw)
        result.setdefault("decision", "continue")
        result.setdefault("next_actions", [])
        return result
    except:
        return {"decision": "continue", "next_actions": []}

@app.post("/generate_ppt")
async def generate_ppt_endpoint(req: GeneratePPTRequest):
    try: llm = get_llm()
    except Exception as e: raise HTTPException(503, str(e))
    
    title = req.title or req.topic.title()
    content_prompt = f"""Generate {req.num_slides} slides about: {req.topic}
Title: {title}
{"Instructions: " + req.additional_instructions if req.additional_instructions else ""}
Return JSON array: hero, two_column, big_number, comparison, quote types.
First=hero, last=quote. Be specific."""
    
    try:
        raw = _call(llm, content_prompt).strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        slides_data = json.loads(raw)
        if not isinstance(slides_data, list): slides_data = [slides_data]
    except:
        slides_data = [
            {"type": "hero", "title": title, "subtitle": "Presentation"},
            {"type": "quote", "quote": "The future belongs to those who prepare.", "author": "Malcolm X"}
        ]
    
    try:
        from ppt_generator_pro import (PREMIUM_THEMES, create_hero_slide, create_two_column_slide,
            create_big_number_slide, create_comparison_slide, create_quote_slide, add_slide_transition_advanced)
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches
        
        prs = PptxPresentation()
        prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
        colors = PREMIUM_THEMES.get(req.theme.lower().replace(" ","_"), PREMIUM_THEMES.get("gamma_modern"))
        
        for idx, sd in enumerate(slides_data):
            t = sd.get("type", "two_column")
            try:
                if t == "hero": s = create_hero_slide(prs, sd.get("title",title), sd.get("subtitle",""), colors)
                elif t == "big_number": s = create_big_number_slide(prs, sd.get("title",""), sd.get("number",""), sd.get("description",""), colors)
                elif t == "comparison": s = create_comparison_slide(prs, sd.get("title",""), sd.get("left_title","A"), sd.get("left_items",[]), sd.get("right_title","B"), sd.get("right_items",[]), colors)
                elif t == "quote": s = create_quote_slide(prs, sd.get("quote",""), sd.get("author",""), colors)
                else: s = create_two_column_slide(prs, sd.get("title",""), sd.get("left_content",[]), sd.get("right_content",[]), colors)
                add_slide_transition_advanced(s, ["zoom","reveal","morph","fade"][idx%4])
            except Exception as se: logger.warning(f"Slide {idx}: {se}")
        
        save_dir = _resolve_save_dir(req.save_path)
        os.makedirs(save_dir, exist_ok=True)
        filename = re.sub(r'[\\/:*?"<>|]', '_', title) + ".pptx"
        full_path = os.path.join(save_dir, filename)
        prs.save(full_path)
        return {"success": True, "path": full_path.replace("\\","/"), "slides_count": len(slides_data),
                "message": f"Created '{title}' ({len(slides_data)} slides) at {full_path}"}
    except ImportError as ie: raise HTTPException(500, f"ppt_generator_pro missing: {ie}")
    except Exception as e:
        logger.error(traceback.format_exc())
        raise HTTPException(500, str(e))

# ─── STARTUP ─────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import argparse
    p = argparse.ArgumentParser()
    p.add_argument("--port", type=int, default=8000)
    p.add_argument("--host", default="0.0.0.0")
    args = p.parse_args()
    logger.info(f"Planner: Groq/{GROQ_MODEL}")
    logger.info(f"Vision : {'CogAgent@'+COGAGENT_URL if COGAGENT_URL else 'Gemini/'+GEMINI_MODEL if GEMINI_API_KEY else 'NONE'}")
    uvicorn.run(app, host=args.host, port=args.port)
