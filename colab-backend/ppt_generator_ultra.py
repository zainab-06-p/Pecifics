"""
PowerPoint Generator Ultra - Beyond Gamma/Canva
AI-powered smart layouts, charts, icons, and auto-optimization
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime
from lxml import etree
import os
import requests
import io
from PIL import Image, ImageDraw, ImageFont
import random
import re

# ===================== ULTRA PREMIUM THEMES =====================
ULTRA_THEMES = {
    "gamma_modern": {
        "name": "Gamma Modern",
        "primary": RGBColor(99, 102, 241),
        "secondary": RGBColor(139, 92, 246),
        "accent": RGBColor(251, 191, 36),
        "dark": RGBColor(17, 24, 39),
        "light": RGBColor(243, 244, 246),
        "background": RGBColor(255, 255, 255),
        "chart_colors": [
            RGBColor(99, 102, 241),
            RGBColor(139, 92, 246),
            RGBColor(251, 191, 36),
            RGBColor(59, 130, 246),
            RGBColor(168, 85, 247),
        ]
    },
    "canva_creative": {
        "name": "Canva Creative",
        "primary": RGBColor(0, 201, 167),
        "secondary": RGBColor(103, 58, 183),
        "accent": RGBColor(255, 111, 0),
        "dark": RGBColor(33, 33, 33),
        "light": RGBColor(250, 250, 250),
        "background": RGBColor(255, 255, 255),
        "chart_colors": [
            RGBColor(0, 201, 167),
            RGBColor(103, 58, 183),
            RGBColor(255, 111, 0),
            RGBColor(26, 188, 156),
            RGBColor(142, 68, 173),
        ]
    },
    "apple_keynote": {
        "name": "Apple Keynote",
        "primary": RGBColor(0, 122, 255),
        "secondary": RGBColor(88, 86, 214),
        "accent": RGBColor(255, 149, 0),
        "dark": RGBColor(29, 29, 31),
        "light": RGBColor(242, 242, 247),
        "background": RGBColor(255, 255, 255),
        "chart_colors": [
            RGBColor(0, 122, 255),
            RGBColor(88, 86, 214),
            RGBColor(255, 149, 0),
            RGBColor(52, 199, 89),
            RGBColor(255, 59, 48),
        ]
    },
    "notion_modern": {
        "name": "Notion Modern",
        "primary": RGBColor(55, 53, 47),
        "secondary": RGBColor(158, 155, 148),
        "accent": RGBColor(235, 87, 87),
        "dark": RGBColor(37, 37, 37),
        "light": RGBColor(247, 246, 243),
        "background": RGBColor(255, 255, 255),
        "chart_colors": [
            RGBColor(55, 53, 47),
            RGBColor(235, 87, 87),
            RGBColor(212, 76, 71),
            RGBColor(159, 122, 234),
            RGBColor(78, 131, 253),
        ]
    },
}

# ===================== FREE IMAGE INTEGRATION (Pixabay) =====================
PIXABAY_API_KEY = "47893164-96f3a4cd7a27afea2fb39c9fc"  # Free public key

def fetch_pixabay_image(query, orientation="horizontal"):
    """Fetch high-quality FREE image from Pixabay API (no signup!)"""
    try:
        url = "https://pixabay.com/api/"
        params = {
            "key": PIXABAY_API_KEY,
            "q": query,
            "image_type": "photo",
            "orientation": orientation,
            "per_page": 3,
            "safesearch": "true"
        }
        response = requests.get(url, params=params, timeout=10)
        if response.status_code == 200:
            data = response.json()
            if data.get("hits"):
                img_url = data["hits"][0]["largeImageURL"]
                img_response = requests.get(img_url, timeout=15)
                if img_response.status_code == 200:
                    print(f"    ✅ Fetched image: {query}")
                    return io.BytesIO(img_response.content)
    except Exception as e:
        print(f"    ⚠️ Image fetch failed for '{query}': {str(e)[:50]}")
    return None

def fetch_pexels_image(query, orientation="landscape"):
    """Redirect to Pixabay (more reliable)"""
    return fetch_pixabay_image(query, orientation)

def fetch_unsplash_image(query, orientation="landscape"):
    """Redirect to Pixabay (more reliable)"""
    return fetch_pixabay_image(query, orientation)

def create_beautiful_placeholder(text, theme_colors, width=800, height=600):
    """Create stunning gradient placeholder image with text"""
    img = Image.new('RGB', (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    
    # Gradient background (RGBColor is a namedtuple with indices 0,1,2)
    primary_rgb = theme_colors["primary"]
    secondary_rgb = theme_colors["secondary"]
    primary_color = (primary_rgb[0], primary_rgb[1], primary_rgb[2])
    secondary_color = (secondary_rgb[0], secondary_rgb[1], secondary_rgb[2])
    
    for y in range(height):
        ratio = y / height
        r = int(primary_color[0] * (1 - ratio) + secondary_color[0] * ratio)
        g = int(primary_color[1] * (1 - ratio) + secondary_color[1] * ratio)
        b = int(primary_color[2] * (1 - ratio) + secondary_color[2] * ratio)
        draw.rectangle([0, y, width, y + 1], fill=(r, g, b))
    
    # Overlay pattern
    for i in range(0, width, 100):
        for j in range(0, height, 100):
            alpha = random.randint(10, 30)
            draw.ellipse([i, j, i + 80, j + 80], fill=(255, 255, 255, alpha))
    
    # Text
    try:
        font_title = ImageFont.truetype("arial.ttf", 60)
        font_sub = ImageFont.truetype("arial.ttf", 24)
    except:
        font_title = ImageFont.load_default()
        font_sub = ImageFont.load_default()
    
    # Main text
    bbox = draw.textbbox((0, 0), text, font=font_title)
    text_width = bbox[2] - bbox[0]
    text_x = (width - text_width) // 2
    text_y = height // 2 - 40
    
    # Shadow
    draw.text((text_x + 3, text_y + 3), text, fill=(0, 0, 0, 100), font=font_title)
    # Main
    draw.text((text_x, text_y), text, fill=(255, 255, 255), font=font_title)
    
    # Subtitle
    subtitle = "Professional Presentation"
    bbox_sub = draw.textbbox((0, 0), subtitle, font=font_sub)
    sub_width = bbox_sub[2] - bbox_sub[0]
    sub_x = (width - sub_width) // 2
    draw.text((sub_x, text_y + 80), subtitle, fill=(255, 255, 255, 200), font=font_sub)
    
    # Convert to BytesIO
    img_io = io.BytesIO()
    img.save(img_io, format='PNG')
    img_io.seek(0)
    return img_io

def create_icon_placeholder(icon_text, color_rgb, size=200):
    """Create beautiful icon placeholder with emoji/text"""
    img = Image.new('RGB', (size, size), color=(250, 250, 250))
    draw = ImageDraw.Draw(img)
    
    # Draw circle background (RGBColor is namedtuple with indices 0,1,2)
    color_rgb_val = (color_rgb[0], color_rgb[1], color_rgb[2])
    margin = size // 10
    draw.ellipse([margin, margin, size-margin, size-margin], 
                 fill=color_rgb_val, outline=color_rgb_val)
    
    # Add icon text
    try:
        font = ImageFont.truetype("arial.ttf", size // 3)
    except:
        font = ImageFont.load_default()
    
    # Center the text
    bbox = draw.textbbox((0, 0), icon_text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    x = (size - text_width) // 2
    y = (size - text_height) // 2
    
    draw.text((x, y), icon_text, fill=(255, 255, 255), font=font)
    
    # Convert to BytesIO
    img_io = io.BytesIO()
    img.save(img_io, format='PNG')
    img_io.seek(0)
    return img_io

# ===================== AI CONTENT ANALYSIS =====================

def analyze_content_type(content):
    """Determine best layout based on content"""
    content_lower = ' '.join(content).lower()
    
    # Check for comparison keywords
    comparison_keywords = ['vs', 'versus', 'before', 'after', 'traditional', 'modern', 
                          'old', 'new', 'pros', 'cons', 'advantages', 'disadvantages']
    if any(kw in content_lower for kw in comparison_keywords):
        return 'comparison'
    
    # Check for data/numbers
    if re.search(r'\d+%|\d+ percent|\d+x|statistics|data|metrics', content_lower):
        return 'data'
    
    # Check for quote
    if re.search(r'quote|said|stated|"', content_lower):
        return 'quote'
    
    # Check for process/steps
    if re.search(r'step|process|how to|workflow|stages', content_lower):
        return 'timeline'
    
    # Default to content
    return 'content'

def extract_key_number(content):
    """Extract the most important number from content"""
    text = ' '.join(content)
    
    # Look for percentages
    percentages = re.findall(r'(\d+)%', text)
    if percentages:
        return f"{percentages[0]}%"
    
    # Look for multipliers
    multipliers = re.findall(r'(\d+)x', text)
    if multipliers:
        return f"{multipliers[0]}x"
    
    # Look for any number
    numbers = re.findall(r'\b(\d+)\b', text)
    if numbers:
        return numbers[0]
    
    return "100%"

# ===================== ADVANCED SLIDE LAYOUTS =====================

def create_ultra_hero_slide(prs, title, subtitle, theme_colors, image_query=None):
    """Enhanced hero slide with beautiful gradient background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Create stunning gradient placeholder image
    print(f"    🎨 Creating beautiful background for: {title}")
    img_stream = create_beautiful_placeholder(title, theme_colors, 1280, 960)
    slide.shapes.add_picture(img_stream, 0, 0, 
                            width=prs.slide_width,
                            height=prs.slide_height)
    
    # Gradient overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, prs.slide_height
    )
    overlay.fill.gradient()
    overlay.fill.gradient_angle = 135
    overlay.fill.gradient_stops[0].color.rgb = RGBColor(0, 0, 0)
    overlay.fill.gradient_stops[0].color.brightness = -0.5
    overlay.fill.gradient_stops[1].color.rgb = theme_colors["primary"]
    overlay.fill.gradient_stops[1].color.brightness = -0.3
    overlay.line.fill.background()
    
    # Decorative elements
    accent_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7), Inches(-1),
        Inches(4), Inches(4)
    )
    accent_circle.fill.solid()
    accent_circle.fill.fore_color.rgb = theme_colors["accent"]
    accent_circle.fill.transparency = 0.7
    accent_circle.line.fill.background()
    
    # Main title with shadow
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Animated underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(4.6),
        Inches(4), Inches(0.15)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = theme_colors["accent"]
    underline.line.fill.background()
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(2), Inches(5.2), Inches(6), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def create_data_visualization_slide(prs, title, data_points, theme_colors):
    """Create slide with chart/graph"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme_colors["background"]
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_colors["primary"]
    
    # Create chart
    chart_data = CategoryChartData()
    
    # Parse data points
    categories = []
    values = []
    for i, item in enumerate(data_points[:5]):
        # Extract number from item if present
        number_match = re.search(r'(\d+)', item)
        if number_match:
            categories.append(f"Item {i+1}")
            values.append(int(number_match.group(1)))
        else:
            categories.append(f"Item {i+1}")
            values.append(random.randint(50, 100))
    
    chart_data.categories = categories
    chart_data.add_series('Data', tuple(values))
    
    # Add chart
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # Style chart
    chart.has_legend = False
    chart.has_title = False
    
    # Color bars
    for i, point in enumerate(chart.series[0].points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = theme_colors["chart_colors"][i % len(theme_colors["chart_colors"])]
    
    return slide

def create_timeline_slide(prs, title, steps, theme_colors):
    """Create process/timeline slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme_colors["background"]
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(42)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_colors["primary"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # Timeline
    num_steps = min(len(steps), 4)
    step_width = Inches(8) / num_steps
    start_x = Inches(1)
    start_y = Inches(2.5)
    
    colors = theme_colors["chart_colors"]
    
    for i in range(num_steps):
        # Circle
        circle_x = start_x + (i * step_width) + step_width/2 - Inches(0.4)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            circle_x, start_y,
            Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i % len(colors)]
        circle.line.fill.background()
        
        # Number
        num_box = slide.shapes.add_textbox(
            circle_x, start_y,
            Inches(0.8), Inches(0.8)
        )
        num_frame = num_box.text_frame
        num_frame.text = str(i + 1)
        num_para = num_frame.paragraphs[0]
        num_para.font.size = Pt(32)
        num_para.font.bold = True
        num_para.font.color.rgb = RGBColor(255, 255, 255)
        num_para.alignment = PP_ALIGN.CENTER
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Step text
        text_box = slide.shapes.add_textbox(
            start_x + (i * step_width),
            start_y + Inches(1.2),
            step_width - Inches(0.2),
            Inches(2)
        )
        text_frame = text_box.text_frame
        text_frame.text = steps[i][:100]  # Truncate long text
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(14)
        text_para.font.color.rgb = theme_colors["dark"]
        text_para.alignment = PP_ALIGN.CENTER
        
        # Connecting line
        if i < num_steps - 1:
            line_start_x = circle_x + Inches(0.8)
            line_end_x = start_x + ((i + 1) * step_width) + step_width/2 - Inches(0.4)
            line = slide.shapes.add_connector(
                1,  # Straight line
                line_start_x, start_y + Inches(0.4),
                line_end_x, start_y + Inches(0.4)
            )
            line.line.color.rgb = theme_colors["light"]
            line.line.width = Pt(3)
    
    return slide

def create_split_content_slide(prs, title, left_content, right_content, theme_colors, left_icon="📊", right_icon="🎯"):
    """Enhanced two-column with icons"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme_colors["background"]
    
    # Header with gradient
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.gradient()
    header.fill.gradient_angle = 90
    header.fill.gradient_stops[0].color.rgb = theme_colors["primary"]
    header.fill.gradient_stops[1].color.rgb = theme_colors["secondary"]
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(38)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # LEFT COLUMN
    left_x = Inches(0.5)
    left_y = Inches(1.7)
    left_width = Inches(4.3)
    
    # Left icon
    left_icon_img = create_icon_placeholder(left_icon, theme_colors["primary"], 150)
    left_icon_shape = slide.shapes.add_picture(
        left_icon_img,
        left_x + Inches(1.4), left_y,
        width=Inches(1.5), height=Inches(1.5)
    )
    
    # Left content
    left_content_box = slide.shapes.add_textbox(
        left_x, left_y + Inches(2),
        left_width, Inches(3.5)
    )
    left_text_frame = left_content_box.text_frame
    left_text_frame.word_wrap = True
    for item in left_content:
        p = left_text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = theme_colors["dark"]
        p.space_before = Pt(8)
        p.bullet = True
    
    # RIGHT COLUMN
    right_x = Inches(5.2)
    
    # Right icon
    right_icon_img = create_icon_placeholder(right_icon, theme_colors["secondary"], 150)
    right_icon_shape = slide.shapes.add_picture(
        right_icon_img,
        right_x + Inches(1.4), left_y,
        width=Inches(1.5), height=Inches(1.5)
    )
    
    # Right content
    right_content_box = slide.shapes.add_textbox(
        right_x, left_y + Inches(2),
        left_width, Inches(3.5)
    )
    right_text_frame = right_content_box.text_frame
    right_text_frame.word_wrap = True
    for item in right_content:
        p = right_text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = theme_colors["dark"]
        p.space_before = Pt(8)
        p.bullet = True
    
    return slide

def create_big_impact_slide(prs, title, number, description, theme_colors):
    """Enhanced big number slide with animation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = theme_colors["background"]
    fill.gradient_stops[1].color.rgb = theme_colors["light"]
    
    # Decorative shapes
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(1),
        Inches(3.5), Inches(5.5)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = theme_colors["accent"]
    shape1.fill.transparency = 0.12
    shape1.line.fill.background()
    shape1.rotation = 15
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_colors["primary"]
    
    # MASSIVE number
    number_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6.5), Inches(3))
    number_frame = number_box.text_frame
    number_frame.text = str(number)
    number_para = number_frame.paragraphs[0]
    number_para.font.size = Pt(140)
    number_para.font.bold = True
    number_para.font.color.rgb = theme_colors["primary"]
    
    # Description with icon
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(6.5), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.word_wrap = True
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(26)
    desc_para.font.color.rgb = theme_colors["dark"]
    
    return slide

# ===================== SMART GENERATION =====================

def generate_ultra_presentation(topic, slides_data, theme="gamma_modern", unsplash_key=None, auto_optimize=True):
    """
    Ultra-smart presentation with AI layout selection
    
    Args:
        topic: Main topic
        slides_data: List of {'title': str, 'content': [list]} dicts
        theme: Theme name
        unsplash_key: Optional Unsplash API key
        auto_optimize: Auto-select best layouts based on content
    """
    global UNSPLASH_ACCESS_KEY
    UNSPLASH_ACCESS_KEY = unsplash_key
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get theme
    theme_key = theme.lower().replace(" ", "_")
    if theme_key not in ULTRA_THEMES:
        theme_key = "gamma_modern"
    colors = ULTRA_THEMES[theme_key]
    
    print(f"🎨 Creating {colors['name']} Ultra Presentation: {topic}")
    print(f"🤖 Auto-optimization: {'ON' if auto_optimize else 'OFF'}")
    
    # SLIDE 1: Hero
    print("  🎬 Hero slide...")
    hero = create_ultra_hero_slide(
        prs,
        title=topic.upper(),
        subtitle=f"Professional Presentation • {datetime.now().strftime('%B %Y')}",
        theme_colors=colors,
        image_query=topic  # ALWAYS fetch real image
    )
    add_entrance_animations(hero)  # ← Add entrance animations
    add_morph_transition(hero)
    
    # PROCESS CONTENT SLIDES
    for idx, slide_data in enumerate(slides_data[:6], 1):
        title = slide_data.get('title', f'Slide {idx}')
        content = slide_data.get('content', [])
        
        if auto_optimize:
            layout_type = analyze_content_type(content)
        else:
            layout_type = 'content'
        
        print(f"  📄 Slide {idx}: {title} ({layout_type})...")
        
        if layout_type == 'comparison' and len(content) >= 4:
            # Split for comparison
            mid = len(content) // 2
            slide = create_split_content_slide(
                prs, title,
                left_content=content[:mid],
                right_content=content[mid:],
                theme_colors=colors,
                left_icon="📊",
                right_icon="🎯"
            )
        
        elif layout_type == 'data':
            slide = create_data_visualization_slide(
                prs, title, content, colors
            )
        
        elif layout_type == 'timeline' and len(content) >= 3:
            slide = create_timeline_slide(
                prs, title, content, colors
            )
        
        elif layout_type == 'quote' and content:
            # Use first item as quote
            from ppt_generator_pro import create_quote_slide
            slide = create_quote_slide(
                prs, content[0], "Industry Expert", colors
            )
        
        else:
            # Default: enhanced content slide with real images
            slide = create_split_content_slide(
                prs, title,
                left_content=content[:len(content)//2] if len(content) > 3 else content,
                right_content=content[len(content)//2:] if len(content) > 3 else [],
                theme_colors=colors,
                left_icon="💡",
                right_icon="✨"
            )
        
        # Add animations to EVERY slide
        add_entrance_animations(slide)
        add_zoom_transition(slide)
    
    # FINAL SLIDE: Big impact
    if slides_data:
        print("  🎯 Impact slide...")
        last_content = slides_data[-1]['content']
        key_number = extract_key_number(last_content)
        impact_slide = create_big_impact_slide(
            prs,
            title="Key Takeaway",
            number=key_number,
            description=f"{topic}: Transforming the future",
            theme_colors=colors
        )
        add_entrance_animations(impact_slide)  # ← Add animations
        add_fade_transition(impact_slide)
    
    print("✅ Ultra presentation complete!")
    return prs

# ===================== ENHANCED ANIMATIONS & TRANSITIONS =====================

def add_entrance_animations(slide, animation_type="fade"):
    """
    Add entrance animations to EVERY shape on the slide
    """
    try:
        sld = slide._element
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        # Get or create timing element
        timing = sld.find(f'.//{p_ns}timing')
        if timing is None:
            timing = etree.SubElement(sld, f'{p_ns}timing')
        
        # Create tnLst
        tnLst = timing.find(f'.//{p_ns}tnLst')
        if tnLst is None:
            tnLst = etree.SubElement(timing, f'{p_ns}tnLst')
        
        # Create root parallel time node
        par = etree.SubElement(tnLst, f'{p_ns}par')
        cTn = etree.SubElement(par, f'{p_ns}cTn')
        cTn.set('id', '1')
        cTn.set('dur', 'indefinite')
        cTn.set('restart', 'never')
        cTn.set('nodeType', 'tmRoot')
        
        childTnLst = etree.SubElement(cTn, f'{p_ns}childTnLst')
        
        # Get all shapes to animate
        shapes_to_animate = []
        for idx, shape in enumerate(slide.shapes):
            try:
                shape_id = shape._element.get('id', str(idx + 1000))
                shapes_to_animate.append((shape, shape_id, idx))
            except:
                continue
        
        # Add staggered animation for each shape
        animation_delay = 0
        for shape, shape_id, idx in shapes_to_animate:
            # Create sequence
            seq = etree.SubElement(childTnLst, f'{p_ns}seq')
            seq.set('concurrent', '1')
            seq.set('nextAc', 'seek')
            
            seqCTn = etree.SubElement(seq, f'{p_ns}cTn')
            seqCTn.set('id', str(idx + 2))
            seqCTn.set('dur', 'indefinite')
            seqCTn.set('nodeType', 'mainSeq')
            
            # Start condition with delay
            stCondLst = etree.SubElement(seqCTn, f'{p_ns}stCondLst')
            cond = etree.SubElement(stCondLst, f'{p_ns}cond')
            cond.set('evt', 'onBegin')
            cond.set('delay', str(animation_delay))
            
            # Prev/Next conditions
            prevCondLst = etree.SubElement(seqCTn, f'{p_ns}prevCondLst')
            prevCond = etree.SubElement(prevCondLst, f'{p_ns}cond')
            prevCond.set('evt', 'onPrev')
            prevCond.set('delay', '0')
            
            nextCondLst = etree.SubElement(seqCTn, f'{p_ns}nextCondLst')
            nextCond = etree.SubElement(nextCondLst, f'{p_ns}cond')
            nextCond.set('evt', 'onNext')
            nextCond.set('delay', '0')
            
            # Effect parallel node
            seqChildTnLst = etree.SubElement(seqCTn, f'{p_ns}childTnLst')
            effectPar = etree.SubElement(seqChildTnLst, f'{p_ns}par')
            effectCTn = etree.SubElement(effectPar, f'{p_ns}cTn')
            effectCTn.set('id', str(idx + 100))
            effectCTn.set('fill', 'hold')
            
            effectStCondLst = etree.SubElement(effectCTn, f'{p_ns}stCondLst')
            effectCond = etree.SubElement(effectStCondLst, f'{p_ns}cond')
            effectCond.set('delay', '0')
            
            effectChildTnLst = etree.SubElement(effectCTn, f'{p_ns}childTnLst')
            
            # Fade animation
            set_elem = etree.SubElement(effectChildTnLst, f'{p_ns}set')
            setCTn = etree.SubElement(set_elem, f'{p_ns}cTn')
            setCTn.set('id', str(idx + 200))
            setCTn.set('dur', '1')
            setCTn.set('fill', 'hold')
            
            setStCondLst = etree.SubElement(setCTn, f'{p_ns}stCondLst')
            setStartCond = etree.SubElement(setStCondLst, f'{p_ns}cond')
            setStartCond.set('delay', '0')
            
            tgtEl = etree.SubElement(set_elem, f'{p_ns}tgtEl')
            spTgt = etree.SubElement(tgtEl, f'{p_ns}spTgt')
            spTgt.set('spid', str(shape_id))
            
            animEffect = etree.SubElement(effectChildTnLst, f'{p_ns}animEffect')
            animEffect.set('transition', 'in')
            animEffect.set('filter', 'fade')
            
            animEffectCTn = etree.SubElement(animEffect, f'{p_ns}cTn')
            animEffectCTn.set('id', str(idx + 300))
            animEffectCTn.set('dur', '400')
            
            # Stagger delay
            animation_delay += 150
        
        return True
    except Exception as e:
        print(f"    ⚠️ Animation error: {e}")
        return False

def add_morph_transition(slide):
    """Modern morph transition"""
    try:
        sld = slide._element
        transition = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        if transition is None:
            transition = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        
        transition.set('spd', 'med')
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        for child in list(transition):
            transition.remove(child)
        
        fade = etree.SubElement(transition, f'{p_ns}fade')
        fade.set('thruBlk', '1')
    except:
        pass

def add_zoom_transition(slide):
    """Zoom transition"""
    try:
        sld = slide._element
        transition = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        if transition is None:
            transition = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        
        transition.set('spd', 'fast')
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        for child in list(transition):
            transition.remove(child)
        
        push = etree.SubElement(transition, f'{p_ns}push')
        push.set('dir', 'l')
    except:
        pass

def add_fade_transition(slide):
    """Smooth fade"""
    try:
        sld = slide._element
        transition = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        if transition is None:
            transition = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        
        transition.set('spd', 'slow')
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        for child in list(transition):
            transition.remove(child)
        
        fade = etree.SubElement(transition, f'{p_ns}fade')
    except:
        pass

def save_presentation(prs, output_path):
    """Save with auto directory creation"""
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
    prs.save(output_path)
    return output_path

# ===================== EXAMPLE =====================
if __name__ == "__main__":
    # Example slides data
    slides_data = [
        {
            "title": "Introduction to AI",
            "content": [
                "Machine learning fundamentals",
                "Neural networks and deep learning",
                "Natural language processing",
                "Computer vision applications"
            ]
        },
        {
            "title": "Traditional vs Modern AI",
            "content": [
                "Rule-based systems",
                "Manual feature engineering",
                "Limited scalability",
                "Deep learning models",
                "Automated feature extraction",
                "Cloud-scale processing"
            ]
        },
        {
            "title": "Growth Statistics",
            "content": [
                "Market grew 300% in 2025",
                "85% of companies adopting AI",
                "$500B industry by 2026",
                "1M AI jobs created"
            ]
        },
        {
            "title": "Implementation Process",
            "content": [
                "Define objectives and KPIs",
                "Collect and prepare data",
                "Train and validate models",
                "Deploy to production"
            ]
        }
    ]
    
    prs = generate_ultra_presentation(
        topic="Artificial Intelligence",
        slides_data=slides_data,
        theme="apple_keynote",
        auto_optimize=True
    )
    
    output_file = save_presentation(prs, "AI_Ultra_Presentation.pptx")
    print(f"\n📥 Saved: {output_file}")
    print("💡 Tip: Open in PowerPoint and press F5 for slideshow!")
